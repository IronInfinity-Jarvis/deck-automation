import streamlit as st
import pandas as pd
import os
import re
from pptx import Presentation as ptx
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.util import Inches, Pt
import pandas as pd
from pptx.oxml.xmlchemy import OxmlElement
from io import BytesIO
from spire.presentation.common import *
from spire.presentation import Presentation as SpirePresentation
from PIL import Image
from tempfile import NamedTemporaryFile
from textwrap import wrap
import numpy as np
from sklearn.preprocessing import MinMaxScaler
import tiktoken
from typing import Optional
from google import generativeai as genai
import google.generativeai as genai
import textwrap 
from dotenv import load_dotenv
import openpyxl

load_dotenv() 
# Function to load data safely
@st.cache_data
def load_sheet_names(data_file):
    if data_file is not None:
        xls = pd.ExcelFile(data_file)
        return xls.sheet_names
    return []

def load_data(data_file, sheet_name):
    df = pd.read_excel(data_file, sheet_name=sheet_name, header=[0, 1])
    df.columns = df.columns.to_frame().ffill(axis=0)
    return df

def process_data(df, change, p2, p1):
    df[("Pen. Change (%)", change)] = df[('Penetration (%)', p2)] - df[('Penetration (%)', p1)]
    df[("Avg. Vol. Growth (%)", change)] = (df["Average Volume (grams/HH)", p2] / df["Average Volume (grams/HH)", p1]) - 1
    df[("Avg. Val. Growth (%)", change)] = (df["Average Value (Rs./HH)", p2] / df["Average Value (Rs./HH)", p1]) - 1
    df[("Avg. NOP. Growth (%)", change)] = (df["Average NOP (nos./HH)", p2] / df["Average NOP (nos./HH)", p1]) - 1
    df[("Avg. Pack Size Growth (%)", change)] = (df["Average Pack Size (grams/Unit)", p2] / df["Average Pack Size (grams/Unit)", p1]) - 1
    df[("Avg.No.of Brands Growth (%)", change)] = (df["Average Brands consumed (nos./HH)", p2] / df["Average Brands consumed (nos./HH)", p1]) - 1
    df[("Avg. PPU Growth (%)", change)] = (df["Average Price Per Unit (Rs./Unit)", p2] / df["Average Price Per Unit (Rs./Unit)", p1]) - 1
    df[("Avg. PPG Growth (%)", change)] = (df["Average Price per gram (Rs./gram)", p2] / df["Average Price per gram (Rs./gram)", p1]) - 1
    df[("Vol. Growth (%)", change)] = (df["Volume ('000 kgs.)", p2] / df["Volume ('000 kgs.)", p1]) - 1
    df[("Val. Growth (%)", change)] = (df["Value (Rs.)", p2] / df["Value (Rs.)", p1]) - 1
    df[("NOP. Growth (%)", change)] = (df["NOP", p2] / df["NOP", p1]) - 1
    df[("SOR Vol. Change (%)", change)] = df["SOR by Volume", p2] - df["SOR by Volume", p1]
    df[("SOR Val. Change (%)", change)] = df["SOR by Value", p2] - df["SOR by Value", p1]
    df[("SOR Packs change (%)", change)] = df["SOR by Packs", p2] - df["SOR by Packs", p1]
    df[("Vol. Share change (%)", change)] = df["Volume Share", p2] - df["Volume Share", p1]
    df[("Val. Share change (%)", change)] = df["Value Share", p2] - df["Value Share", p1]
    df[("Pack Share change (%)", change)] = df["Pack Share", p2] - df["Pack Share", p1]
    df[("Proj. HH Growth (%)", change)] = (df["Projected HH", p2] / df["Projected HH", p1]) - 1
    return df
# AI Header Generation Agent
from typing import Tuple, Dict

class AIDataInterpreter:
    def __init__(self):
        self.client = genai.configure(api_key=os.getenv('GEMINI_API_KEY'))

    def _structure_data(self, df: pd.DataFrame, time_period: str, 
                       change_period: Optional[str] = None) -> str:
        """Process data WITHOUT changing column names"""
        structured = []
        
        for _, row in df.iterrows():
            brand = row.iloc[0]
            metrics = []
            
            for (metric, period), value in row.items():
                if period == 'Time Period':
                    continue
                
                # Format based on original metric names
                fmt_value = self._format_value(metric, value)
                metrics.append(f"{metric}: {fmt_value}")

            structured.append(f"Brand: {brand}\n" + "\n".join(metrics))
        
        time_context = f"Current Period: {time_period}"
        if change_period:
            time_context += f"\nComparison Period: {change_period}"
        
        return f"{time_context}\n\n" + "\n\n".join(structured)

    def _format_value(self, metric: str, value) -> str:
        """Format values using ORIGINAL metric names"""
        if any(kw in metric for kw in ['Penetration', 'Share', 'Growth', 'Change', 'SOR']):
            return f"{value*100:.1f}%" if isinstance(value, (float, int)) else str(value)
            
        if any(kw in metric for kw in ['Rs.', 'Price', 'PPU', 'PPG']):
            return f"₹{value:,.2f}" if isinstance(value, (float, int)) else str(value)
            
        if any(kw in metric for kw in ['Volume', 'Pack Size', 'NOP', 'HH']):
            return f"{value:,.0f}" if isinstance(value, (float, int)) else str(value)
            
        return str(value)

    def generate_insights(self, df: pd.DataFrame, time_period: str,
                         change_period: Optional[str] = None) -> str:
        """Generate insights using original column names"""
        structured_data = self._structure_data(df, time_period, change_period)
        prompt = f"""Analyze the given FMCG data using the original metric names exactly as they appear. Extract key insights from {structured_data}
        focusing on below points if present:
        Penetration trends across brands and categories
        Shifts in value and volume and pack share over time not actual value,volume or projected households
        Price per gram (PPG) variations and their implications
        Pack size developments and consumer preferences
        SOR metric which speaks about consumer loyalty
        consumer in table represent usual consumer used brands other than professional name written brands.
        hair serum, mask, shampoo, conditioner are subcategories under hair care if these appear they are not brands there are categories similarly for face and body.
        Brand and category-level insights, including top brands across metrics and emerging small brands
        Your insights should be comprehensive, data-driven, and out-of-the-box, combining multiple metrics to uncover meaningful trends. Write in proper paragraph format, ensuring concise yet complete insights that fit well within a PowerPoint slide created using the python-pptx library with, dont repeat what i gave you in prompt!:
        Slide width: 13.33 inches
        Slide height: 7.5 inches
        Use shorter sentences and multiple lines to optimize readability without exceeding slide height.
        Example insights for reference (do not copy):
        >50% of households adopt post-wash treatments in hair care; affluent consumers spend ~₹1,900 annually on shampoo and post-wash products.
        Shampoo and conditioner growth is driven by increased consumption, whereas hair serum growth stems from more consumers trying the category.
        New-age brands like Mamaearth are rapidly gaining market share, while L'Oréal Professionnel dominates salon-exclusive brands, reaching 1 in 4 households.
        Kerastase Paris leads in volume growth, while professional brands such as Mamaearth expand—except Biolage and OGX, which face declines.
        Each insight should be at the same level of depth and clarity as the examples above. Ensure a strong narrative flow and full-sentence structure, avoiding fragmented insights."""
        models = genai.GenerativeModel('gemini-2.0-flash')
        response = models.generate_content(prompt)
        raw_text = response.text.strip()
        return self._format_for_ppt(raw_text)
    
    def _format_for_ppt(self, text: str) -> str:
        """Format text for PowerPoint constraints"""
        # Split into logical sections first
        sections = text.split(' | ')
        
        formatted_lines = []
        current_line = []
        current_length = 0
        
        # Adjust based on your font testing (characters per line)
        max_line_length = 100  # Start with 100 chars for 12pt font
        
        for section in sections:
            if current_length + len(section) > max_line_length and current_line:
                formatted_lines.append(" | ".join(current_line))
                current_line = []
                current_length = 0
                
            current_line.append(section)
            current_length += len(section) + 3  # Account for " | " separator
            
        # Add any remaining content
        if current_line:
            formatted_lines.append(" | ".join(current_line))
        
        # Handle long individual insights with word wrap
        final_lines = []
        for line in formatted_lines:
            if len(line) > max_line_length:
                wrapped = textwrap.fill(line, width=max_line_length, 
                                    break_long_words=False)
                final_lines.extend(wrapped.split('\n'))
            else:
                final_lines.append(line)
        return '\n'.join(final_lines)

def generate_tag_patterns(tag_list):
    tag_patterns = []
    for tag in tag_list:
        # This pattern splits based on space, hyphen, underscore, or camelCase
        split_tag = re.split(r'[\s_/*.+-]+', tag)
        patterns = r'.*'.join(map(re.escape, split_tag))  # Escape special characters, join with '.*'
        tag_patterns.append(patterns)
    return tag_patterns

# Convert PowerPoint slide to image
def slide_to_image(slide):
    image_stream = slide.SaveAsImage()  # This returns a Stream object
    img_bytes = image_stream.ToArray()  # Convert the Stream to bytes
    # Use PIL to open the image from the byte data
    img = Image.open(BytesIO(img_bytes))
    # Optionally, you can dispose of the image object to free resources
    image_stream.Dispose()
    return img

# Function to convert PowerPoint slides to images
def convert_ppt_to_images(ppt_buffer):
    # Write the PPT buffer to a temporary PPTX file
    with NamedTemporaryFile(delete=False, suffix=".pptx") as tmp_file:
        tmp_file.write(ppt_buffer.getvalue())
        tmp_file_path = tmp_file.name
    # Load the PPTX file with Spire Presentation
    ppt = SpirePresentation()
    ppt.LoadFromFile(tmp_file_path)  # Load the PowerPoint file
    images = []
    # Convert each slide to an image
    for slide in ppt.Slides:
        img = slide_to_image(slide)  # Convert slide to image
        images.append(img)
    # Clean up the temporary PowerPoint file
    os.remove(tmp_file_path)
    return images

# Display the slides in a slideshow format
def display_slideshow(images):
    for img in images:
        container_width = 500  # Set the width you want for the container
        img = img.resize((container_width, int(container_width * img.height / img.width)))
        st.image(img, use_container_width=True)  # Display each slide as an image in Streamlit

def extract_table_data(table):
    """Extracts data from PowerPoint table and converts it to a Pandas DataFrame."""
    data = []
    for row in table.rows:
        row_data = [cell.text for cell in row.cells]  # Extract text from each cell
        data.append(row_data)
    
    # Convert to DataFrame
    df = pd.DataFrame(data)
    return df

def slide_table_formatter_loreal_2(table, change_lst, no_cols, tag_first_indices, loreal_matching_indices):
    for col_idx in range(no_cols):
        for row_idx in range(len(table.rows)):
            row = table.rows[row_idx]  # Access the specific row
            cell = row.cells[col_idx]  # Get the cell in the specified column
            text_frame = cell.text_frame
            text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
            cell.fill.solid()  # Fill the cell with a solid color
            cell.fill.fore_color.rgb = RGBColor(255, 255, 255)  # White color
            if row_idx == 0:
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.alignment = PP_ALIGN.CENTER
                    for run in paragraph.runs:
                        run.font.size = Pt(15)  # Set font size
                        run.font.bold = False  # Set font to bold
                        run.font.color.rgb = RGBColor(127, 127, 127)
                        run.font.name = 'LOREAL Essentielle'
            elif row_idx == 1:
                for paragraph in cell.text_frame.paragraphs:
                    if row_idx == 3:
                        paragraph.alignment = PP_ALIGN.CENTER
                    for run in paragraph.runs:
                        run.font.size = Pt(13)  # Set font size
                        if row_idx == 1:
                            run.font.bold = False  # Set font to bold
                        else:
                            run.font.bold = True  # Set font to bold
                        run.font.color.rgb = RGBColor(255, 255, 255)
                        run.font.name = 'LOREAL Essentielle'
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(166, 114, 49)  # Blue color
            elif row_idx == 2:
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.alignment = PP_ALIGN.CENTER
                    for run in paragraph.runs:
                        run.font.size = Pt(13)  # Set font size
                        run.font.bold = False  # Set font to bold
                        run.font.color.rgb = RGBColor(127, 127, 127)  # Set font color to gray
                        run.font.name = 'LOREAL Essentielle'
            elif row_idx in tag_first_indices:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(249, 203, 123)  # Blue color
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.alignment = PP_ALIGN.CENTER
                    for run in paragraph.runs:
                        run.font.size = Pt(12)  # Set font size
                        run.font.bold = True  # Set font to not bold
                        run.font.color.rgb = RGBColor(0, 0, 0)  # Set font color to black
            elif row_idx in loreal_matching_indices:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(255, 204, 204)  # Blue color
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.alignment = PP_ALIGN.CENTER
                    for run in paragraph.runs:
                        run.font.size = Pt(12)  # Set font size
                        run.font.bold = True  # Set font to not bold
                        run.font.color.rgb = RGBColor(0, 0, 0)  # Set font color to black
            else:
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.alignment = PP_ALIGN.CENTER
                    for run in paragraph.runs:
                        run.font.size = Pt(12)  # Set font size
                        run.font.bold = True  # Set font to not bold
                        run.font.color.rgb = RGBColor(0, 0, 0)  # Set font color to black
            
            if table.cell(2, col_idx).text in change_lst and row_idx not in [0, 1, 2] and col_idx != 0 and cell.text != '-':
                cell_value = float(cell.text.replace('%', '').strip())
                if cell_value < 0:
                    for paragraph in cell.text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.font.color.rgb = RGBColor(255, 0, 0)
                elif cell_value > 0:
                    for paragraph in cell.text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.font.color.rgb = RGBColor(0, 176, 80)

def slide_table_formatter_loreal_1(table, change_lst, no_cols, loreal_matching_indices):
    for col_idx in range(no_cols):
        for row_idx in range(len(table.rows)):
            row = table.rows[row_idx]  # Access the specific row
            cell = row.cells[col_idx]  # Get the cell in the specified column
            text_frame = cell.text_frame
            text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
            cell.fill.solid()  # Fill the cell with a solid color
            cell.fill.fore_color.rgb = RGBColor(255, 255, 255)  # White color
            if row_idx == 0:
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.alignment = PP_ALIGN.CENTER
                    for run in paragraph.runs:
                        run.font.size = Pt(15)  # Set font size
                        run.font.bold = False  # Set font to bold
                        run.font.color.rgb = RGBColor(127, 127, 127)
                        run.font.name = 'LOREAL Essentielle'
            elif row_idx == 1 or row_idx == 3:
                for paragraph in cell.text_frame.paragraphs:
                    if row_idx == 3:
                        paragraph.alignment = PP_ALIGN.CENTER
                    for run in paragraph.runs:
                        run.font.size = Pt(13)  # Set font size
                        if row_idx == 1:
                            run.font.bold = False  # Set font to bold
                        else:
                            run.font.bold = True  # Set font to bold
                        run.font.color.rgb = RGBColor(255, 255, 255)
                        run.font.name = 'LOREAL Essentielle'
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(166, 114, 49)  # Blue color
            elif row_idx == 2:
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.alignment = PP_ALIGN.CENTER
                    for run in paragraph.runs:
                        run.font.size = Pt(13)  # Set font size
                        run.font.bold = False  # Set font to bold
                        run.font.color.rgb = RGBColor(127, 127, 127)  # Set font color to gray
                        run.font.name = 'LOREAL Essentielle'
            elif row_idx in loreal_matching_indices:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(255, 204, 204)  # Blue color
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.alignment = PP_ALIGN.CENTER
                    for run in paragraph.runs:
                        run.font.size = Pt(12)  # Set font size
                        run.font.bold = True  # Set font to not bold
                        run.font.color.rgb = RGBColor(0, 0, 0)  # Set font color to black
            else:
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.alignment = PP_ALIGN.CENTER
                    for run in paragraph.runs:
                        run.font.size = Pt(12)  # Set font size
                        run.font.bold = True  # Set font to not bold
                        run.font.color.rgb = RGBColor(0, 0, 0)  # Set font color to black
            if table.cell(2, col_idx).text in change_lst and row_idx not in [0, 1, 2] and col_idx != 0 and cell.text != '-':
                cell_value = float(cell.text.replace('%', '').strip())
                if cell_value < 0:
                    for paragraph in cell.text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.font.color.rgb = RGBColor(255, 0, 0)
                elif cell_value > 0:
                    for paragraph in cell.text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.font.color.rgb = RGBColor(0, 176, 80)


def values_inp(metric_change_dict, metric_dict, table, data, start_row, max_rows, no_cols, change_lst, tp, p1 = None, p2 = None):
    keys_list = list(metric_dict.keys())
    # places metric names in first row
    for col_idx in range(no_cols):
        cell = table.cell(2, col_idx+1)
        cell.text = metric_change_dict[keys_list[col_idx]]
    # places names from brand file in first column
    for row_idx in range(start_row, max_rows):
        cell = table.cell(row_idx, 0)
        value = data[('Metric', 'Time Period')].iloc[row_idx - start_row]
        cell.text = value
    # places values in subsequent rows and columns
    for col_idx in range(no_cols):
        for row_idx in range(start_row, max_rows):
            cell = table.cell(row_idx, col_idx+1)
            if keys_list[col_idx] in change_lst:
                value = data[(keys_list[col_idx], f'{p1} vs {p2}')].iloc[row_idx - start_row]
            else:
                value = data[(keys_list[col_idx], tp)].iloc[row_idx - start_row]
            # Set the cell text to the value formatted to 1 decimal place
            percent_lst = ['Penetration(%)', 'Vol. Share(%)', 'Val Share(%)', 'SOR Val.(%)', 'SOR Vol.(%)', 'SOR Pack.(%)', "Pen. Change (%)",
                           'Pack Share(%)', "Pen. Change (%)", "Avg. Vol. Growth (%)", "Avg. Val. Growth (%)", "Avg. NOP. Growth (%)",
                           "Avg. Pack Size Growth (%)", "Avg.No.of Brands Growth (%)", "Avg. PPU Growth (%)", "Avg. PPG Growth (%)", "Vol. Growth (%)",
                           "Val. Growth (%)", "NOP. Growth (%)","SOR Vol. Change (%)", "SOR Val. Change (%)", "SOR Packs change (%)", "Vol. Share change (%)",
                           "Val. Share change (%)", "Pack Share change (%)", 'Proj. HH Growth (%)']

            zero_dec_lst = ["Avg. Vol.(mL)", "Avg. Val.(Rs.)", "Avg. Pack Size(mL)", "Proj. HH", "Val. (Rs.)",
                            "Vol. (Tons)", "NOP (Units)"]
            one_dec_lst = ["Avg. NOP (Units)", "Avg.No.of Brands", "Avg. PPU (Rs./Unit)", "Avg. PPG (Rs./Gram)"]

            if table.cell(2, col_idx+1).text in percent_lst:
                cell.text = f"{value * 100:.1f}%"
            elif table.cell(2, col_idx+1).text in zero_dec_lst:
                cell.text = f"{value:,.0f}"
            elif table.cell(2, col_idx+1).text in one_dec_lst:
                cell.text = f"{value:.1f}"
            cleaned_text = re.sub(r'[^\d.-]', '', cell.text).strip()
            if cleaned_text.lower() in ['nan', 'inf', 'infinity', '']:
                cell.text = '-'


def SubElement(parent, tagname, **kwargs):
    element = OxmlElement(tagname)
    element.attrib.update(kwargs)
    parent.append(element)
    return element


def _set_cell_border(cell, border_color="000000", border_width='6350'):
    """ Hack function to enable the setting of border width and border color
        - left border
        - right border
        - top border
        - bottom border
    """
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()

    # Left Cell Border
    lnL = SubElement(tcPr, 'a:lnL', w=border_width, cap='flat', cmpd='sng', algn='ctr')
    lnL_solidFill = SubElement(lnL, 'a:solidFill')
    lnL_srgbClr = SubElement(lnL_solidFill, 'a:srgbClr', val=border_color)
    lnL_prstDash = SubElement(lnL, 'a:prstDash', val='dashDot')
    lnL_round_ = SubElement(lnL, 'a:round')
    lnL_headEnd = SubElement(lnL, 'a:headEnd', type='none', w='med', len='med')
    lnL_tailEnd = SubElement(lnL, 'a:tailEnd', type='none', w='med', len='med')

    # Right Cell Border
    lnR = SubElement(tcPr, 'a:lnR', w=border_width, cap='flat', cmpd='sng', algn='ctr')
    lnR_solidFill = SubElement(lnR, 'a:solidFill')
    lnR_srgbClr = SubElement(lnR_solidFill, 'a:srgbClr', val=border_color)
    lnR_prstDash = SubElement(lnR, 'a:prstDash', val='dashDot')
    lnR_round_ = SubElement(lnR, 'a:round')
    lnR_headEnd = SubElement(lnR, 'a:headEnd', type='none', w='med', len='med')
    lnR_tailEnd = SubElement(lnR, 'a:tailEnd', type='none', w='med', len='med')

    # Top Cell Border
    lnT = SubElement(tcPr, 'a:lnT', w=border_width, cap='flat', cmpd='sng', algn='ctr')
    lnT_solidFill = SubElement(lnT, 'a:solidFill')
    lnT_srgbClr = SubElement(lnT_solidFill, 'a:srgbClr', val=border_color)
    lnT_prstDash = SubElement(lnT, 'a:prstDash', val='dashDot')
    lnT_round_ = SubElement(lnT, 'a:round')
    lnT_headEnd = SubElement(lnT, 'a:headEnd', type='none', w='med', len='med')
    lnT_tailEnd = SubElement(lnT, 'a:tailEnd', type='none', w='med', len='med')

    # Bottom Cell Border
    lnB = SubElement(tcPr, 'a:lnB', w=border_width, cap='flat', cmpd='sng', algn='ctr')
    lnB_solidFill = SubElement(lnB, 'a:solidFill')
    lnB_srgbClr = SubElement(lnB_solidFill, 'a:srgbClr', val=border_color)
    lnB_prstDash = SubElement(lnB, 'a:prstDash', val='dashDot')
    lnB_round_ = SubElement(lnB, 'a:round')
    lnB_headEnd = SubElement(lnB, 'a:headEnd', type='none', w='med', len='med')
    lnB_tailEnd = SubElement(lnB, 'a:tailEnd', type='none', w='med', len='med')


def slide_gen(pptx, data, no_rows, no_cols, Time_period, change=False, tp1=None, tp2=None, households=None):
    slide = pptx.slides.add_slide(blank_slide_layout)
    left = Inches(0.5)
    top = Inches(0.2)
    width = Inches(13)
    height = Inches(1.5)
    # creating textBox
    txBox = slide.shapes.add_textbox(left, top, width, height)
    # creating textFrames
    tf = txBox.text_frame
    # adding Paragraphs
    p = tf.add_paragraph()
    # adding text
    p.text = "Check Second slide for Insights!"
    p.font.bold = True
    p.font.size = Pt(24)
    p.font.name = 'LOREAL Essentielle'
    x, y, cx, cy = Inches(0.75), Inches(1.5), Inches(12), Inches(6)
    # x = Inches(2) means the table is positioned 2 inches from the left edge of the slide.
    # y = Inches(2) means the table is positioned 2 inches from the top edge of the slide.
    # cx = Inches(4) means the table will be 4 inches wide.
    # cy = Inches(1.5) means the table will be 1.5 inches tall.
    shape = slide.shapes.add_table(no_rows, no_cols, x, y, cx, cy)
    table = shape.table
    for cell in table.iter_cells():
        _set_cell_border(cell)
    first_row = table.rows[0].cells
    first_row[0].merge(first_row[len(table.rows[0].cells) - 1])

    if change:
        if tp1 is None or tp2 is None:
            raise ValueError("tp1 and tp2 must be provided when 'change' is True.")
        first_row[0].text = f"CONSUMPTION  BEHAVIOUR, {tp1} vs {tp2}"
    else:
        first_row[0].text = f"CONSUMPTION  BEHAVIOUR, {Time_period}"

    second_row = table.rows[1].cells
    second_row[0].merge(second_row[len(table.rows[1].cells) - 1])
    second_row[0].text = f"Universe HHs -  {households:,.0f} HHs"
    return table, data

def slide_gen_insights(pptx,out_text):
    slide = pptx.slides.add_slide(blank_slide_layout)
    left = Inches(0.5)
    top = Inches(0.2)
    width = Inches(13)
    height = Inches(1.5)
    # creating textBox
    txBox = slide.shapes.add_textbox(left, top, width, height)
    # creating textFrames
    tf = txBox.text_frame
    # adding Paragraphs
    p = tf.add_paragraph()
    # adding text
    p.text = out_text
    p.font.bold = True
    p.font.size = Pt(12)
    p.font.name = 'LOREAL Essentielle'

if __name__ == "__main__":
    metric_dic_change = {'Penetration (%)': 'Penetration(%)',
                        'Pen. Change (%)': 'Pen. Change (%)',
                        'Projected HH': 'Proj. HH',
                        'Proj. HH Growth (%)': 'Proj. HH Growth (%)',
                        'Value (Rs.)': 'Val. (Rs.)',
                        'Val. Growth (%)': 'Val. Growth (%)',
                        'Average Value (Rs./HH)': 'Avg. Val.(Rs.)',
                        'Avg. Val. Growth (%)': 'Avg. Val. Growth (%)',
                        'Value Share': 'Val Share(%)',
                        'Val. Share change (%)': 'Val. Share change (%)',
                        'SOR by Value': 'SOR Val.(%)',
                        'SOR Val. Change (%)': 'SOR Val. Change (%)',
                        "Volume ('000 kgs.)": 'Vol. (Tons)',
                        'Vol. Growth (%)': 'Vol. Growth (%)',
                        'Average Volume (grams/HH)': 'Avg. Vol.(mL)',
                        'Avg. Vol. Growth (%)': 'Avg. Vol. Growth (%)',
                        'Volume Share': 'Vol. Share(%)',
                        'Vol. Share change (%)': 'Vol. Share change (%)',
                        'SOR by Volume': 'SOR Vol.(%)',
                        'SOR Vol. Change (%)': 'SOR Vol. Change (%)',
                        'NOP': 'NOP (Units)',
                        'NOP. Growth (%)': 'NOP. Growth (%)',
                        'Average NOP (nos./HH)': 'Avg. NOP (Units)',
                        'Avg. NOP. Growth (%)': 'Avg. NOP. Growth (%)',
                        'Pack Share': 'Pack Share(%)',
                        'Pack Share change (%)': 'Pack Share change (%)',
                        'SOR by Packs': 'SOR Pack.(%)',
                        'SOR Packs change (%)': 'SOR Packs change (%)',
                        'Average Price Per Unit (Rs./Unit)': 'Avg. PPU (Rs./Unit)',
                        'Avg. PPU Growth (%)': 'Avg. PPU Growth (%)',
                        'Average Pack Size (grams/Unit)': 'Avg. Pack Size(mL)',
                        'Avg. Pack Size Growth (%)': 'Avg. Pack Size Growth (%)',
                        'Average Price per gram (Rs./gram)': 'Avg. PPG (Rs./Gram)',
                        'Avg. PPG Growth (%)': 'Avg. PPG Growth (%)',
                        'Average Brands consumed (nos./HH)': 'Avg.No.of Brands',
                        'Avg.No.of Brands Growth (%)': 'Avg.No.of Brands Growth (%)'}
    
    metric_mapping = {'Penetration (%)': 'Pen. Change (%)',
                    'Volume Share': 'Vol. Share change (%)',
                    'Average Volume (grams/HH)': 'Avg. Vol. Growth (%)',
                    'Value Share': 'Val. Share change (%)',
                    'Average Value (Rs./HH)': 'Avg. Val. Growth (%)',
                    'Average Pack Size (grams/Unit)': 'Avg. Pack Size Growth (%)',
                    'Average NOP (nos./HH)': 'Avg. NOP. Growth (%)',
                    'Average Brands consumed (nos./HH)': 'Avg.No.of Brands Growth (%)',
                    'Pack Share': 'Pack Share change (%)',
                    'SOR by Value': 'SOR Val. Change (%)',
                    'SOR by Volume': 'SOR Vol. Change (%)',
                    'SOR by Packs': 'SOR Packs change (%)',
                    'Average Price Per Unit (Rs./Unit)': 'Avg. PPU Growth (%)',
                    'Average Price per gram (Rs./gram)': 'Avg. PPG Growth (%)',
                    "Volume ('000 kgs.)": 'Vol. Growth (%)',
                    'NOP': 'NOP. Growth (%)', 'Value (Rs.)': 'Val. Growth (%)',
                    'Projected HH': 'Proj. HH Growth (%)'}
    
    change_growth_list = ["Pen. Change (%)", "Avg. Vol. Growth (%)", "Avg. Val. Growth (%)", "Avg. NOP. Growth (%)",
                          "Avg. Pack Size Growth (%)", "Avg.No.of Brands Growth (%)", "Avg. PPU Growth (%)",
                          "Avg. PPG Growth (%)", "Vol. Growth (%)", "Val. Growth (%)", "NOP. Growth (%)",
                          "SOR Vol. Change (%)", "SOR Val. Change (%)", "SOR Packs change (%)", "Vol. Share change (%)",
                          "Val. Share change (%)", "Pack Share change (%)", 'Proj. HH Growth (%)']
    st.title("PPT Generator 🤖")
    uploaded_file = st.file_uploader("Upload a KPI report Excel file", type=["xlsx", "xls"])
    if uploaded_file is not None:
        sheet_names = load_sheet_names(uploaded_file)
        if sheet_names:
            city = st.selectbox("Please choose the city or sheet name:", sheet_names)
            df = load_data(uploaded_file, city)
            if df is not None:
                st.write("### Preview of Uploaded Data:")
                st.dataframe(df)
                time_periods = [col[1] for col in df.columns if isinstance(col, tuple)]
                metric_keys = [col[0] for col in df.columns if isinstance(col, tuple)]
                unique_tp = list(set(time_periods))
                metric_keys = list(set(metric_keys))
                col1, col2 = st.columns(2)
                if 'Time Period' in unique_tp:
                    unique_tp.remove('Time Period')
                st.title("Select Time Period")
                prefs = ['no','yes']
                type_of_slide = ['Usual_KPI_Table_Loreal','PRO_CON_KPI_Table_Loreal']
                with col1:
                    tp = st.selectbox("Choose Singular Time period for specific period KPI Tables:", unique_tp)
                growth_pref = st.radio("Do you want to see change or growth of metrics as headers in table?", prefs, index = 0)
                if growth_pref == 'yes':
                    with col2:
                        p1 = st.selectbox("Choose Second Timeperiod for growths and changes:", unique_tp, index=1)
                p2 = tp
                pensort_pref = st.radio("Do you want your data to be penetration sorted?", prefs, index = 0)
                filtered_dic = {key: metric_dic_change[key] for key in metric_keys if key in metric_dic_change}
                if growth_pref == 'yes':
                    updated_dic = {metric_mapping[key]: metric_dic_change[metric_mapping[key]] for key in filtered_dic if key in metric_mapping}
                    for k, v in updated_dic.items():
                        filtered_dic.setdefault(k, v)  # Only adds if `k` is not already in `filtered_dic`
                st.title("Select Metrics and Order")
                if "selection_order" not in st.session_state:
                    st.session_state.selection_order = {}
                num_cols = 5  # Number of columns per row
                cols = st.columns(num_cols)
                filtered_dic_sorted = {key : filtered_dic[key] for key in metric_dic_change if key in filtered_dic}
                keys = list(filtered_dic_sorted.keys())
                for idx, key in enumerate(keys):
                    display_name = filtered_dic_sorted[key]
                    col = cols[idx % num_cols]  # Assign column dynamically
                    with col:
                        selected = st.checkbox(f"{display_name}", key=f"check_{key}")
                        if selected:
                            if key not in st.session_state.selection_order:
                                max_order = max(st.session_state.selection_order.values(), default=0) + 1
                                st.session_state.selection_order[key] = max_order
                        else:
                            if key in st.session_state.selection_order:
                                del st.session_state.selection_order[key]  # Remove if unchecked

                # **Reorder the dictionary so that numbers are sequential (1, 2, 3...)**
                sorted_keys = sorted(st.session_state.selection_order, key=lambda k: st.session_state.selection_order[k])
                st.session_state.selection_order = {key: i+1 for i, key in enumerate(sorted_keys)}
                st.write("### Selected Metrics and Their Order")
                st.write(st.session_state.selection_order)

                slide_type = st.selectbox("What kind of slide do you want to generate?", type_of_slide)
                if growth_pref == 'yes':
                    st.write(f"Selected Time Periods: {p1} and {p2}")
                    change = f'{p1} vs {p2}'
                    df = process_data(df, change, p2, p1)
                household = df[('Projected HH', tp)].iloc[0]
                if pensort_pref=="yes":
                    df_sorted = df.sort_values(by=('Penetration (%)', tp), ascending=False)
                    data_in = df_sorted
                else:
                    data_in = df.drop(df.index[0])

                if slide_type == 'PRO_CON_KPI_Table_Loreal':
                    text_list = data_in[('Metric', 'Time Period')].tolist()
                    unique_keywords = set()
                    for text in text_list:
                        words = re.split(r"[-\s]+", text)  # Split by "-" and spaces
                        words = [re.sub(r"[^a-zA-Z0-9]", "", word).lower() for word in words]  # Remove special chars & lowercase
                        unique_keywords.update(words)  # Add to set
                    unique_keywords = list(sorted(unique_keywords))
                    user_tag_list  = st.multiselect("Select the Tag elements:",unique_keywords)
                    num_tag_items = st.number_input("How many items for each tag?", min_value=1)
                    name_lst = data_in[('Metric', 'Time Period')].tolist()
                    df_concat = pd.DataFrame(columns=data_in.columns.to_list()).astype(data_in.dtypes.to_dict())
                    tag_pattern = generate_tag_patterns(user_tag_list)
                    df_list = []
                    first_names = []
                    for pattern in tag_pattern:
                        filtered_words = [word for word in name_lst if re.search(pattern, word, re.IGNORECASE)]
                        df_tag = data_in[data_in[('Metric', 'Time Period')].isin(filtered_words)]
                        df_tag = df_tag.head(num_tag_items)
                        first_names.append(df_tag[('Metric', 'Time Period')].iloc[0])
                        df_list.append(df_tag)
                    df_concat = pd.concat(df_list, ignore_index=True) if df_list else pd.DataFrame(columns=data_in.columns.to_list())
                    matching_indices = df_concat[df_concat[('Metric', 'Time Period')].isin(first_names)].index
                    matching_indices_list = matching_indices.tolist()
                    matching_indices_list = [num+3 for num in matching_indices_list]
                    

                if pensort_pref == 'no':
                    if slide_type == 'PRO_CON_KPI_Table_Loreal':
                        brand_file = st.multiselect("Which elements do you want to see KPI table?", df_concat[('Metric', 'Time Period')].tolist())
                    else:
                        brand_file = st.multiselect("Which elements do you want to see KPI table?", data_in[('Metric', 'Time Period')].tolist())
                elif pensort_pref == "yes":
                    if slide_type != 'PRO_CON_KPI_Table_Loreal':
                        num_of_rows = st.number_input("How many rows of data would you like to see in KPI table?", min_value=1)

                if slide_type == 'Usual_KPI_Table_Loreal':
                    if pensort_pref == 'no':
                        filtered_data = data_in[data_in[('Metric', 'Time Period')].isin(brand_file)]
                        no_rows = len(brand_file) + 3
                    elif pensort_pref == 'yes':
                        filtered_data = data_in
                        no_rows = num_of_rows + 3

                elif slide_type == 'PRO_CON_KPI_Table_Loreal':
                    if pensort_pref == 'no':
                        filtered_data = df_concat[df_concat[('Metric', 'Time Period')].isin(brand_file)]
                        no_rows = len(brand_file) + 3
                    elif pensort_pref == "yes":
                        filtered_data = df_concat
                        no_rows = num_tag_items*len(user_tag_list)+3
                max_rows_avail = len(filtered_data[('Metric', 'Time Period')].to_list())
                if no_rows-3 > max_rows_avail:
                    st.warning(f"selected more no. of rows than actually present in filtered data, hence setting no of rows to {max_rows_avail}")
                    no_rows = max_rows_avail + 3      
                filtered_data = filtered_data.iloc[:no_rows-3]  
                selected_metrics = [key for key in st.session_state.selection_order.keys()]
                filtered_columns = []
                for metric in selected_metrics:
                    if metric in change_growth_list:
                        col = (metric, f'{p1} vs {p2}')
                    else:
                        col = (metric, tp)
                    filtered_columns.append(col)
                entity_column = filtered_data.columns[0]
                filtered_columns = [entity_column] + filtered_columns
                filtered_data = filtered_data[filtered_columns]
                filtered_data = filtered_data.reset_index(drop=True)
                st.write("### Preview of Data filtered by your selections!")
                st.dataframe(filtered_data)
                pink_higlights = ["L'oreal","kerastase","Biolage","matrix"]
                pink_tag_pattern = generate_tag_patterns(pink_higlights)
                loreal_filtered_words_lst = []
                loreal_name_lst = filtered_data[('Metric', 'Time Period')].tolist()
                for pattern in pink_tag_pattern:
                    loreal_filtered_words_lst = loreal_filtered_words_lst + [word for word in loreal_name_lst if re.search(pattern, word, re.IGNORECASE)]
                loreal_matching_indices = filtered_data[filtered_data[('Metric', 'Time Period')].isin(loreal_filtered_words_lst)].index.tolist()
                loreal_matching_indices = [num+3 for num in loreal_matching_indices]

                if st.button("Process Data"):
                    final_dict = st.session_state.selection_order
                    no_cols = max(final_dict.values())
                    ppt = ptx()
                    ppt.slide_width = Inches(13.33)  # Width
                    ppt.slide_height = Inches(7.5)  # Height
                    blank_slide_layout = ppt.slide_layouts[6]
                    table, data_in2 = slide_gen(ppt, filtered_data, no_rows, no_cols+1, tp, change=True if growth_pref=='yes' else False, tp1=p1 if growth_pref=='yes' else None, tp2=p2 if growth_pref=='yes' else None, households=household)
                    values_inp(metric_dic_change, final_dict, table, data_in2, 3, no_rows, no_cols, change_growth_list, tp, p1 = p1 if growth_pref=='yes' else None, p2 = p2 if growth_pref=='yes' else None)
                    if slide_type == 'Usual_KPI_Table_Loreal':
                        slide_table_formatter_loreal_1(table, change_growth_list, no_cols+1, loreal_matching_indices)
                    elif slide_type == 'PRO_CON_KPI_Table_Loreal':
                        slide_table_formatter_loreal_2(table, change_growth_list, no_cols+1, matching_indices_list, loreal_matching_indices)
                    ai_analyst = AIDataInterpreter()
                    header_text = ai_analyst.generate_insights(
                        data_in2,
                        tp,
                        change if growth_pref == 'yes' else None
                    )
                    st.title("Insights derived from your selections!")
                    st.write(header_text)
                    slide_gen_insights(ppt,header_text)
                    df_table = extract_table_data(table)
                    st.title("PowerPoint Table Preview")
                    st.dataframe(df_table)
                    ppt_buffer = BytesIO()
                    ppt.save(ppt_buffer)
                    ppt_buffer.seek(0)
                    # images = convert_ppt_to_images(ppt_buffer)
                    # st.title("PowerPoint Table Formatted")
                    # display_slideshow(images)
                    st.download_button(
                        label="Download PowerPoint",
                        data=ppt_buffer,
                        file_name="generated_presentation.pptx",
                        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                    )
            else:
                st.error("Error: The file could not be read. Please check the format.")
